from flask import Flask, render_template, request, jsonify
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
from nltk.probability import FreqDist
from string import punctuation
import heapq
import os
from PyPDF2 import PdfReader
from pptx import Presentation

app = Flask(__name__)

# NLTK gerekli dosyaları indir
nltk.download('punkt')
nltk.download('stopwords')

def extract_text_from_pdf(file_stream):
    reader = PdfReader(file_stream)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_pptx(file_stream):
    prs = Presentation(file_stream)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def summarize_text(text, num_sentences=3):
    # Metni cümlelere ayır
    sentences = sent_tokenize(text)
    
    # Eğer metin çok kısaysa, direkt döndür
    if len(sentences) <= num_sentences:
        return text
    
    # Stopwords'leri yükle
    stop_words = set(stopwords.words('english'))
    
    # Kelimeleri tokenize et ve stopwords'leri çıkar
    word_tokens = word_tokenize(text.lower())
    word_tokens = [word for word in word_tokens if word not in stop_words and word not in punctuation]
    
    # Kelime frekanslarını hesapla
    freq_dist = FreqDist(word_tokens)
    
    # Cümle skorlarını hesapla
    sentence_scores = {}
    for i, sentence in enumerate(sentences):
        for word in word_tokenize(sentence.lower()):
            if word in freq_dist:
                if i not in sentence_scores:
                    sentence_scores[i] = freq_dist[word]
                else:
                    sentence_scores[i] += freq_dist[word]
    
    # En yüksek skorlu cümleleri seç
    summary_sentences = heapq.nlargest(num_sentences, sentence_scores, key=sentence_scores.get)
    summary_sentences = sorted(summary_sentences)
    
    # Özeti oluştur
    summary = ' '.join([sentences[i] for i in summary_sentences])
    return summary

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/summarize', methods=['POST'])
def summarize():
    if 'file' in request.files:
        file = request.files['file']
        if file.filename != '':
            filename = file.filename.lower()
            if filename.endswith('.txt'):
                text = file.read().decode('utf-8')
            elif filename.endswith('.pdf'):
                text = extract_text_from_pdf(file)
            elif filename.endswith('.pptx'):
                text = extract_text_from_pptx(file)
            else:
                return jsonify({'error': 'Sadece .txt, .pdf veya .pptx dosyaları destekleniyor.'}), 400
        else:
            return jsonify({'error': 'Dosya seçilmedi'}), 400
    else:
        text = request.form.get('text', '')
        if not text:
            return jsonify({'error': 'Metin girilmedi'}), 400
    
    try:
        summary = summarize_text(text)
        return jsonify({'summary': summary})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True) 